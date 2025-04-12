import streamlit as st
import pickle
from pathlib import Path
import streamlit_authenticator as stauth
import parser
from parser import Status
import json
from pptx import Presentation
from io import BytesIO
from streamlit_star_rating import st_star_rating
from data import *

st.set_page_config(page_title="Competitive Programming At University of Haifa", page_icon=":shark:", layout="wide")


di = {}
file_path = Path(__file__).parent / 'hashed_pw.pkl'
try:
    with open(file_path, 'rb') as f:
        di = pickle.load(f)
except Exception as e:
    pass
    # st.error(e)
if di == {}:
    di['all'] = {'password': 'fklkjdlsk', 'name': 'all'}
    di = {"usernames": di}
authenticator = stauth.Authenticate(di, 'cpwebsite', '12345', 3)

if 'reg' not in st.session_state:
    st.session_state['reg'] = 1

if st.session_state['reg'] == 1:
    try:
        if st.button("Register"):
            st.session_state['reg'] = 0
            st.rerun()
        authenticator.login('main', clear_on_submit=False)
        if st.session_state['authentication_status']:
            st.session_state['reg'] = 2
            with open(file_path, 'wb') as f:
                pickle.dump(di, f)
            st.rerun()
    except Exception as e:
        st.error(e)
elif st.session_state['reg'] == 0:
    try:
        if st.button("Login"):
            st.session_state['reg'] = 1
            st.rerun()
        email, \
            username, \
            name = authenticator.register_user(password_hint=False)
        # fields={'First name':'CSES Username', 'Last name': 'CSES Handle (Go to your profile, its the numbers in the URL!)'})
        if email and username and name:
            st.success('User registered successfully')
            st.session_state['reg'] = 3
            st.session_state['authentication_status'] = True
            st.session_state['username'] = username
            di['usernames'][username]['score'] = 0
            with open(file_path, 'wb') as f:
                pickle.dump(di, f)

            st.rerun()
    except Exception as e:
        st.error(e)
elif st.session_state['reg'] == 3:
    with st.form('cses_info'):
        cses_username = st.text_input('CSES Username', key='cses_username')
        di['usernames'][st.session_state.get('username')]['cses_username'] = cses_username
        cses_handle = st.text_input('CSES Handle (Go to your profile, it\'s the numbers in the URL!)',
                                    key='cses_handle')
        di['usernames'][st.session_state.get('username')]['cses_handle'] = cses_handle
        etgar_num = st.text_input('What is the number of your ETGAR? (18/19/20)',
                                    key='etgar')
        di['usernames'][st.session_state.get('username')]['etgar'] = etgar_num

        with open(file_path, 'wb') as f:
            pickle.dump(di, f)
        submitted = st.form_submit_button('Submit')
        if cses_username and cses_handle and submitted and etgar_num:
            st.session_state['reg'] = 2
            st.rerun()


def week(list_of_questions, list_of_locked, stars, tasks, totoff):
    global di
    su = 0
    l = list_of_questions
    l2 = list_of_locked
    lc, mc, rc = st.columns(3)
    stars_size = 20
    with lc:
        k = [st.link_button(f"Problem {totoff + i + 1}", f'https://cses.fi/problemset/task/{tid}') for i, tid in enumerate(l)]
        st.text('Finish these problems to unlock more challenging ones!')
        # print(tasks)
    with mc:
        k = [st_star_rating("", 5,stars[i], stars_size, read_only=True,dark_theme=True, key=f's{totoff+i}') for i, tid in enumerate(l)]
    with rc:
        st.write('Here you can see the status of your problems:')
        k = [st.badge(f'Problem {totoff + i + 1}', icon=":material/check:", color="green") if tasks[tid] == Status.AC
             else st.badge(f'Problem {totoff + i + 1}', color='gray') if tasks[tid] == Status.NAT else
        st.badge(f'Problem {totoff + i + 1}', icon=":material/close:", color="red")
             for i, tid in enumerate(l)]
        p = [tasks[tid] for tid in l]
        for i in range(len(l)):
            di['usernames'][st.session_state.get('username')][str(totoff + i)] = p[i]
        su = p.count(Status.AC)
        with open(file_path, 'wb') as f:
            pickle.dump(di, f)
    if su == len(p):
        st.subheader('More challenging problems unlocked!')
        lc, mc, rc = st.columns(3)
        off = len(l)
        with lc:
            k = [st.link_button(f"Problem {totoff + i + 1 + off}", f'https://cses.fi/problemset/task/{tid}') for i, tid in
                 enumerate(l2)]
        with mc:
            k = [st_star_rating("", 5, stars[i+off], stars_size, read_only=True, dark_theme=True,
                                key=f's{totoff + i+off}') for i, tid in enumerate(l2)]
        with rc:
            # st.write('Here you can mark the problems you have completed (they will be saved on your next visit):')
            k2 = [st.badge(f'Problem {totoff + i + off + 1}', icon=":material/check:", color="green") if tasks[tid] == Status.AC
                  else st.badge(f'Problem {totoff + i + off + 1}', color='gray') if tasks[tid] == Status.NAT else
            st.badge(f'Problem {totoff + i + off + 1}', icon=":material/close:", color="red")
                  for i, tid in enumerate(l2)]
            p = [tasks[tid] for tid in l2]
            for i in range(len(l2)):
                di['usernames'][st.session_state.get('username')][str(totoff + i + off)] = p[i]
            su += p.count(Status.AC)
            with open(file_path, 'wb') as f:
                pickle.dump(di, f)
        if su == off + len(p):
            st.success('Congrats! That is all for this week!')
        st.subheader(f'So far you have completed {su}/{off + len(p)} problems this week!')
    else:
        st.subheader(f'So far you have completed {su}/{len(p)} problems this week!')
    return totoff + len(stars)


if st.session_state.get('authentication_status') and st.session_state.get('reg') == 2:
    with st.container():
        authenticator.logout('Logout', 'sidebar')


    def homepage():
        with st.container():
            st.title("Competitive Programming At University of Haifa")
            st.write("Welcome to the Competitive Programming At University of Haifa website!")
            st.write("This website is designed to help students learn and practice competitive programming.")

        with st.container():
            if not st.session_state.get('authentication_status'):
                st.rerun()
            cses_handle = di['usernames'][st.session_state.get('username')].get('cses_handle')
            # print(cses_handle)
            tasks = parser.get_user_info(cses_handle)
            st.write("---")
            st.header("Week One - Introduction to Competitive Programming")
            st.write("Here is the presentation for this week:")
            pr = Presentation('presentations/week1.pptx')
            bo = BytesIO()
            pr.save(bo)
            st.download_button(label='Week 1 Presentation', data=bo.getvalue(), file_name='Competitive-Programming-week-1.pptx')
            st.write("""
                   This week we will be introducing the basics of competitive programming.
                   Here are some problems to get you started:
                   """)
            new_off = week(week1u, week1l, week1s, tasks, 0)

        with st.container():
            if not st.session_state.get('authentication_status'):
                st.rerun()
            cses_handle = di['usernames'][st.session_state.get('username')].get('cses_handle')
            # print(cses_handle)
            tasks = parser.get_user_info(cses_handle)
            st.write("---")
            st.header("Week Two - Basic Data Structures")
            st.write("Tomorrow a presentation will be available...")
            lc, rc = st.columns(2)
            with lc:
                st.write("Here is the presentation for this week:")
                pr = Presentation('presentations/week2-19.pptx') if di['usernames'][st.session_state.get('username')].get('etgar') == '19' \
                    else Presentation('presentations/week2-18.pptx')
                bo = BytesIO()
                pr.save(bo)
                st.download_button(label='Week 2 Presentation', data=bo.getvalue(), file_name='Competitive-Programming-week-2.pptx')
            with rc:
                st.write("Here is the recording from this week (notice that it may be partial during to technical difficulties):")
                st.page_link('https://drive.google.com/file/d/16UVeVeJU4kw3UeJgp_52zh8SKfWt5oC-/view?usp=sharing',
                             label='Recording Drive Link', icon='🇬🇬')
            st.write("""
                   This week we attempt to teach you about data structures that are used in c++ and competitive programming.
                   Here are some problems to get you started:
                   """)
            new_off = week(week2u, week2l, week2s, tasks, new_off)
    pg = st.navigation([homepage, 'leaderboard.py', 'profile.py'])
    pg.run()

